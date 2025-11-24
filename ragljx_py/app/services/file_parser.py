"""文件解析服务"""
import io
import os
from typing import Optional
from loguru import logger

# 文档解析库
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    from striprtf.striprtf import rtf_to_text
except ImportError:
    rtf_to_text = None

import json
import csv


class FileParserService:
    """文件解析服务"""
    
    SUPPORTED_EXTENSIONS = {
        '.txt', '.md', '.pdf', '.docx', '.xlsx', 
        '.pptx', '.html', '.htm', '.csv', '.json', 
        '.xml', '.rtf'
    }
    
    def __init__(self):
        pass
    
    def is_supported(self, filename: str) -> bool:
        """检查文件是否支持解析"""
        ext = os.path.splitext(filename)[1].lower()
        return ext in self.SUPPORTED_EXTENSIONS
    
    def parse(self, file_content: bytes, mime_type: str, filename: str = "") -> tuple[bool, str, str]:
        """
        解析文件内容
        
        Args:
            file_content: 文件二进制内容
            mime_type: MIME 类型
            filename: 文件名
        
        Returns:
            (success, content, error_message)
        """
        try:
            # 根据文件扩展名或 MIME 类型选择解析器
            ext = os.path.splitext(filename)[1].lower() if filename else ""
            
            if ext in ['.txt', '.md'] or 'text/plain' in mime_type or 'text/markdown' in mime_type:
                content = self._parse_text(file_content)
            elif ext == '.pdf' or 'application/pdf' in mime_type:
                content = self._parse_pdf(file_content)
            elif ext == '.docx' or 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' in mime_type:
                content = self._parse_docx(file_content)
            elif ext == '.xlsx' or 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' in mime_type:
                content = self._parse_xlsx(file_content)
            elif ext == '.pptx' or 'application/vnd.openxmlformats-officedocument.presentationml.presentation' in mime_type:
                content = self._parse_pptx(file_content)
            elif ext in ['.html', '.htm'] or 'text/html' in mime_type:
                content = self._parse_html(file_content)
            elif ext == '.csv' or 'text/csv' in mime_type:
                content = self._parse_csv(file_content)
            elif ext == '.json' or 'application/json' in mime_type:
                content = self._parse_json(file_content)
            elif ext == '.xml' or 'application/xml' in mime_type or 'text/xml' in mime_type:
                content = self._parse_xml(file_content)
            elif ext == '.rtf' or 'application/rtf' in mime_type:
                content = self._parse_rtf(file_content)
            else:
                return False, "", f"Unsupported file type: {ext or mime_type}"
            
            if not content or not content.strip():
                return False, "", "No content extracted from file"
            
            return True, content, ""
        
        except Exception as e:
            logger.error(f"Error parsing file: {e}")
            return False, "", str(e)
    
    def _parse_text(self, content: bytes) -> str:
        """解析文本文件"""
        try:
            return content.decode('utf-8')
        except UnicodeDecodeError:
            # 尝试其他编码
            for encoding in ['gbk', 'gb2312', 'latin-1']:
                try:
                    return content.decode(encoding)
                except:
                    continue
            raise ValueError("Unable to decode text file")
    
    def _parse_pdf(self, content: bytes) -> str:
        """解析 PDF 文件"""
        if fitz is None:
            raise ImportError("PyMuPDF not installed")
        
        text_parts = []
        with fitz.open(stream=content, filetype="pdf") as doc:
            for page in doc:
                text_parts.append(page.get_text())
        
        return "\n\n".join(text_parts)
    
    def _parse_docx(self, content: bytes) -> str:
        """解析 DOCX 文件"""
        if Document is None:
            raise ImportError("python-docx not installed")
        
        doc = Document(io.BytesIO(content))
        text_parts = [para.text for para in doc.paragraphs if para.text.strip()]
        return "\n\n".join(text_parts)
    
    def _parse_xlsx(self, content: bytes) -> str:
        """解析 XLSX 文件"""
        if load_workbook is None:
            raise ImportError("openpyxl not installed")
        
        wb = load_workbook(io.BytesIO(content), read_only=True)
        text_parts = []
        
        for sheet in wb.worksheets:
            text_parts.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    text_parts.append(row_text)
        
        return "\n".join(text_parts)
    
    def _parse_pptx(self, content: bytes) -> str:
        """解析 PPTX 文件"""
        if Presentation is None:
            raise ImportError("python-pptx not installed")
        
        prs = Presentation(io.BytesIO(content))
        text_parts = []
        
        for i, slide in enumerate(prs.slides, 1):
            text_parts.append(f"Slide {i}:")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        
        return "\n\n".join(text_parts)
    
    def _parse_html(self, content: bytes) -> str:
        """解析 HTML 文件"""
        if BeautifulSoup is None:
            raise ImportError("beautifulsoup4 not installed")
        
        html = content.decode('utf-8', errors='ignore')
        soup = BeautifulSoup(html, 'html.parser')
        
        # 移除 script 和 style 标签
        for script in soup(["script", "style"]):
            script.decompose()
        
        return soup.get_text(separator="\n", strip=True)
    
    def _parse_csv(self, content: bytes) -> str:
        """解析 CSV 文件"""
        text = content.decode('utf-8', errors='ignore')
        reader = csv.reader(io.StringIO(text))
        rows = ["\t".join(row) for row in reader]
        return "\n".join(rows)
    
    def _parse_json(self, content: bytes) -> str:
        """解析 JSON 文件"""
        data = json.loads(content.decode('utf-8'))
        return json.dumps(data, indent=2, ensure_ascii=False)
    
    def _parse_xml(self, content: bytes) -> str:
        """解析 XML 文件"""
        if BeautifulSoup is None:
            raise ImportError("beautifulsoup4 not installed")
        
        xml = content.decode('utf-8', errors='ignore')
        soup = BeautifulSoup(xml, 'xml')
        return soup.get_text(separator="\n", strip=True)
    
    def _parse_rtf(self, content: bytes) -> str:
        """解析 RTF 文件"""
        if rtf_to_text is None:
            raise ImportError("striprtf not installed")
        
        rtf = content.decode('utf-8', errors='ignore')
        return rtf_to_text(rtf)

